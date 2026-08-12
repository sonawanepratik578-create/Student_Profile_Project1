# Create a PowerPoint presentation about IC 741 with diagrams using python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "IC 741 Operational Amplifier"
slide.placeholders[1].text = "Symbol, Pin Diagram, Working and Applications"

# Slide 2: Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Introduction to IC 741"
tf = slide.shapes.placeholders[1].text = (
    "• IC 741 is a very popular Operational Amplifier (Op-Amp).\n"
    "• It is used for signal amplification in analog circuits.\n"
    "• Developed for general purpose amplification.\n"
    "• Works with dual power supply.\n"
    "• Used in filters, comparators, integrators and oscillators."
)

# Slide 3: Op-Amp Symbol
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Operational Amplifier Symbol"

left = Inches(3)
top = Inches(2)
width = Inches(3)
height = Inches(2)

shape = slide.shapes.add_shape(1, left, top, width, height)  # triangle-like placeholder
shape.text = "741\nOp-Amp"

slide.shapes.add_textbox(Inches(2.2), Inches(2.5), Inches(1), Inches(0.5)).text = "+"
slide.shapes.add_textbox(Inches(2.2), Inches(3.1), Inches(1), Inches(0.5)).text = "-"
slide.shapes.add_textbox(Inches(6.2), Inches(2.8), Inches(1), Inches(0.5)).text = "Output"

# Slide 4: Pin Diagram
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "IC 741 Pin Diagram"

box = slide.shapes.add_shape(1, Inches(3.5), Inches(2), Inches(2), Inches(3))
box.text = "IC 741"

pins_left = [
    "1 Offset Null",
    "2 Inverting Input (-)",
    "3 Non-Inverting Input (+)",
    "4 V- (Negative Supply)"
]

pins_right = [
    "5 Offset Null",
    "6 Output",
    "7 V+ (Positive Supply)",
    "8 NC (No Connection)"
]

for i, p in enumerate(pins_left):
    slide.shapes.add_textbox(Inches(2), Inches(2 + i*0.6), Inches(1.5), Inches(0.5)).text = p

for i, p in enumerate(pins_right):
    slide.shapes.add_textbox(Inches(5.6), Inches(2 + i*0.6), Inches(2), Inches(0.5)).text = p

# Slide 5: Working Principle
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Working Principle"
slide.placeholders[1].text = (
    "• IC 741 amplifies the difference between two input signals.\n"
    "• It has two inputs: Inverting (-) and Non-Inverting (+).\n"
    "• Output = Gain × (V+ − V−).\n"
    "• If signal is applied to inverting input, output is inverted.\n"
    "• If applied to non-inverting input, output remains in phase."
)

# Slide 6: Important Characteristics
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Important Characteristics"
slide.placeholders[1].text = (
    "• High voltage gain (~100,000).\n"
    "• High input impedance.\n"
    "• Low output impedance.\n"
    "• Wide range of applications.\n"
    "• Requires external components for many circuits."
)

# Slide 7: Applications
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Applications of IC 741"
slide.placeholders[1].text = (
    "• Audio amplifiers\n"
    "• Active filters\n"
    "• Voltage followers\n"
    "• Integrators and differentiators\n"
    "• Comparators\n"
    "• Oscillator circuits"
)

# Slide 8: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
slide.placeholders[1].text = (
    "IC 741 is one of the most widely used operational amplifiers. "
    "It is easy to use, reliable, and forms the basis of many analog "
    "electronic circuits used in engineering and practical applications."
)

file_path = "/mnt/data/IC741_Operational_Amplifier_Presentation.pptx"
prs.save(file_path)

file_path